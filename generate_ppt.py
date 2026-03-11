from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_bullet_slide(title_text, bullets, image_path=None, split=False):
    layout_idx = 1 if image_path is None else 5
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title = slide.shapes.title
    title.text = title_text
    
    if image_path:
        # add image
        # add text box below or side
        if split:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
            try:
                pic = slide.shapes.add_picture(image_path, Inches(5), Inches(2), width=Inches(4.5))
            except Exception as e:
                print(f"Could not add image {image_path}: {e}")
        else:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
            try:
                pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(3.0), height=Inches(4.0))
            except Exception as e:
                print(f"Could not add image {image_path}: {e}")
        
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            if bullet.startswith("  -"):
                p.text = bullet.strip(" -")
                p.level = 1
                p.font.size = Pt(16)
            else:
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
    else:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            if bullet.startswith("  -"):
                p.text = bullet.strip(" -")
                p.level = 1
            else:
                p.text = bullet
                p.level = 0

add_title_slide(
    "Towards Zero-Waste Inference\nHardware-Efficient SNNs",
    "ECE274 Neuromorphic Computing Project\nvia Dynamic Gatekeeping & Early Exit"
)

add_bullet_slide(
    "1. Motivation: The Edge Computing Bottleneck",
    [
        "The Promise of SNNs: Event-driven sparsity offers tremendous energy advantages for always-on edge sensors.",
        "The Reality of CMOS Hardware: Fetching data from Memory (SRAM) costs 25x more energy than compute (MAC).",
        "  - SRAM Read (45nm): ~ 5 pJ per access.",
        "  - INT8 Addition: ~ 0.2 pJ per access.",
        "The Problem: Current SNN research focuses on algorithmic accuracy but ignores the hardware cost of redundant memory reads caused by sensor noise."
    ]
)

add_bullet_slide(
    "2. Our Hardware-First Solution",
    [
        "We propose a natively digital, co-designed architecture that intercepts and drops redundant spikes before they can trigger expensive SRAM reads.",
        "Three Core Contributions:",
        "  - 1. Dynamic Gatekeeper Filtering: Pre-processing input noise and burst redundancy.",
        "  - 2. Adaptive Activity Regulation: Hardware-efficient threshold adaptation to silence hidden-layer 'spike storms'.",
        "  - 3. Early-Exit FSM: Terminating the temporal processing window the moment the network reaches a confident prediction."
    ]
)

add_bullet_slide(
    "3. System-Level Architecture: Multi-Tile NoC",
    [
        "To avoid a massive centralized memory bottleneck, our architecture is partitioned into Tiles.",
        "Each layer operates as an independent hardware engine with strictly local, private SRAM."
    ],
    "Hardware_Architecture/multi_tile_noc_diagram.png"
)

add_bullet_slide(
    "4. Inside the Datapath: The SNN Hardware Tile",
    [
        "Each Tile contains everything it needs to execute its layer:",
        "  - Local SRAM Weight Bank: Stores INT8 quantized weights.",
        "  - Sparse MAC Array: Adds weights conditionally based on binary spikes.",
        "  - LIF Neuron Array: Handles leaky integration and thresholding."
    ],
    "Hardware_Architecture/snn_tile_diagram.png", split=True
)

add_bullet_slide(
    "5. Core Innovation A: Dynamic Gatekeeper",
    [
        "Sits in front of the Conv1 tile to pre-filter raw camera input.",
        "  - Importance Monitor: Filters Spatial Noise using an array of 4-bit saturating counters with global decay.",
        "  - Burst Redundancy Filter: Filters Temporal Noise dropping consecutive identical spikes."
    ],
    "circuit_dynamic_gatekeeper.png"
)

add_bullet_slide(
    "6. Core Innovation B: Early-Exit FSM",
    [
        "Sits at the output classifier layer to truncate execution in the temporal dimension.",
        "Concept: Easy inputs (like a clear digit '1') don't need all 20 timesteps to be recognized.",
        "Hardware: 10 simple integer accumulators track class spikes. When any hits ConfThreshold=8, a global Freeze signal is raised.",
        "Result: The entire pipeline stops, saving 100% of the energy for the remaining temporal window."
    ]
)

add_bullet_slide(
    "7. Software Model & Training (PyTorch)",
    [
        "Hardware optimizations must be learned! We modeled the exact digital mechanics directly into a custom PyTorch LeNet-5 SNN.",
        "Encoding: N-MNIST DVS data is cropped to 28x28 and binned to T=20 discrete steps.",
        "Learning: Surrogate Gradient descent (STBP) with a cosine annealing schedule.",
        "Co-Design: Gatekeeper logic, adaptive thresholds, and early exit are active during the forward pass."
    ]
)

add_bullet_slide(
    "8. Results: Massive SRAM Reduction",
    [
        "By combining spatial filtering and temporal truncation, average SRAM reads dropped by 81.6%."
    ],
    "fig1_cumulative_sram.png"
)

add_bullet_slide(
    "9. Results: Energy & Latency Scaling",
    [
        "Because SRAM dominates power, this 81.6% read reduction maps directly to an 81.6% reduction in estimated inference energy.",
        "Average latency improved by 1.8x due to Early Exit stopping at T=11.4 on average."
    ],
    "fig5_latency_speedup.png"
)

add_bullet_slide(
    "10. Conclusion",
    [
        "Hardware-First SNNs: We proved that optimizing for algorithmic sparsity is not enough; we must optimize for memory access.",
        "Zero-Waste Inference: Simple, highly-efficient digital blocks dramatically reduce the power footprint.",
        "Future Work: Deploying the Verilog RTL onto an FPGA to gather physical synthesis, power, and timing reports."
    ]
)

prs.save("ECE274_Presentation.pptx")
print("Presentation generated successfully!")
